"""
agents/ingestion/office_parser.py
------------------------------------
DOCX and PPTX ingestion.

Provenance note: a .docx is a flowing document -- a paragraph doesn't live
at a fixed pixel position the way a PDF block does, so `bbox` is None and
`locator` (e.g. "paragraph_12", "table_0_row_3") carries provenance instead.
A .pptx slide, by contrast, IS a fixed canvas -- every shape has a real
(left, top, width, height) in EMUs, so we convert that to points and give
PPTX blocks a genuine bbox, same coordinate convention as the PDF parser.
"""

from __future__ import annotations

from docx import Document as DocxDocument
from pptx import Presentation

from core.schemas import ContentBlock, SourceFormat
from agents.ingestion.base import next_block_id

EMU_PER_POINT = 12700  # python-pptx reports shape geometry in EMUs


def parse_docx(path: str, source_file: str | None = None) -> tuple[list[ContentBlock], int, list[str]]:
    source_file = source_file or str(path)
    warnings: list[str] = []
    doc = DocxDocument(path)
    blocks: list[ContentBlock] = []
    counter = 0

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        counter += 1
        blocks.append(
            ContentBlock(
                block_id=next_block_id("para", counter),
                source_file=source_file,
                source_format=SourceFormat.DOCX,
                page=None,
                locator=f"paragraph_{i}",
                text=text,
                bbox=None,
                confidence=1.0,
                extraction_method="docx_paragraph",
            )
        )

    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if not row_text.strip(" |"):
                continue
            counter += 1
            blocks.append(
                ContentBlock(
                    block_id=next_block_id("tbl", counter),
                    source_file=source_file,
                    source_format=SourceFormat.DOCX,
                    page=None,
                    locator=f"table_{t_idx}_row_{r_idx}",
                    text=row_text,
                    bbox=None,
                    confidence=1.0,
                    extraction_method="docx_table_row",
                )
            )

    if not blocks:
        warnings.append("DOCX contained no extractable paragraph or table text.")

    # .docx has no fixed page count without a rendering engine (page breaks
    # are a layout-time concept, not stored in the XML). We surface "None"
    # honestly rather than guessing -- this is logged so a future sprint can
    # decide whether to add a real page count via LibreOffice headless.
    return blocks, None, warnings  # type: ignore[return-value]


def parse_pptx(path: str, source_file: str | None = None) -> tuple[list[ContentBlock], int, list[str]]:
    source_file = source_file or str(path)
    warnings: list[str] = []
    prs = Presentation(path)
    blocks: list[ContentBlock] = []
    counter = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            counter += 1

            bbox = None
            if (
                shape.left is not None
                and shape.top is not None
                and shape.width is not None
                and shape.height is not None
            ):
                x0 = shape.left / EMU_PER_POINT
                y0 = shape.top / EMU_PER_POINT
                x1 = (shape.left + shape.width) / EMU_PER_POINT
                y1 = (shape.top + shape.height) / EMU_PER_POINT
                bbox = (x0, y0, x1, y1)

            blocks.append(
                ContentBlock(
                    block_id=next_block_id(f"slide{slide_idx}", counter),
                    source_file=source_file,
                    source_format=SourceFormat.PPTX,
                    page=slide_idx,  # slide number doubles as "page" for UI consistency
                    locator=f"slide_{slide_idx}_shape_{shape_idx}",
                    text=text,
                    bbox=bbox,
                    confidence=1.0,
                    extraction_method="pptx_shape_text",
                )
            )

    page_count = len(prs.slides)
    if not blocks:
        warnings.append("PPTX contained no extractable shape text.")

    return blocks, page_count, warnings
