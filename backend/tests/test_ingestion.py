"""
tests/test_ingestion.py
--------------------------
Sprint 1 exit criteria: every supported format produces ContentBlocks with
correct provenance, the OCR fallback path actually triggers on scanned
pages, long documents (20+ pages) are handled, and hyperlink ingestion
works -- tested against a local HTTP server rather than the live internet,
since unit tests shouldn't depend on external network availability.
"""

from __future__ import annotations

import functools
import http.server
import threading
import time

import pytest

from agents.ingestion import office_parser, pdf_parser
from agents.ingestion.base import UnsupportedFormatError
from agents.ingestion.dispatcher import ingest, ingest_many
from core.schemas import SourceFormat


# ---------------------------------------------------------------------------
# Fixtures: build tiny synthetic documents directly with the same libraries
# scripts/generate_samples.py uses, so these tests are self-contained and
# don't depend on samples/ existing or being up to date.
# ---------------------------------------------------------------------------

@pytest.fixture
def digital_pdf(tmp_path):
    import fitz
    path = tmp_path / "digital.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Policy Number: AUTO-2026-00981", fontsize=12)
    page.insert_text((50, 80), "Date of Loss: 2026-06-12", fontsize=12)
    page2 = doc.new_page()
    page2.insert_text((50, 50), "Repair Estimate: $4250.00", fontsize=12)
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def scanned_pdf(tmp_path):
    """A PDF page with an embedded image and NO text layer -- must trigger OCR."""
    import io
    import fitz
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (800, 300), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 24)
    except OSError:
        font = ImageFont.load_default()
    draw.text((30, 30), "Policy Number: PROP-2026-55210", fill="black", font=font)
    draw.text((30, 80), "Cause of Loss: Wind damage", fill="black", font=font)

    buf = io.BytesIO()
    img.save(buf, format="PNG")

    path = tmp_path / "scanned.pdf"
    doc = fitz.open()
    page = doc.new_page(width=800, height=300)
    page.insert_image(fitz.Rect(0, 0, 800, 300), stream=buf.getvalue())
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def long_pdf(tmp_path):
    import fitz
    path = tmp_path / "long.pdf"
    doc = fitz.open()
    n_pages = 22
    for i in range(1, n_pages + 1):
        page = doc.new_page()
        page.insert_text((50, 50), f"Page {i} of {n_pages}", fontsize=12)
    doc.save(str(path))
    doc.close()
    return path, n_pages


@pytest.fixture
def sample_docx(tmp_path):
    from docx import Document
    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Patient Name: Daniel Osei")
    doc.add_paragraph("Member ID: MBR-77231908")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Total Billed Amount"
    table.rows[0].cells[1].text = "$9,840.00"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Claim Overview"
    slide.placeholders[1].text_frame.text = "Policyholder: Priya Nair"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_image(tmp_path):
    from PIL import Image, ImageDraw, ImageFont
    path = tmp_path / "photo.png"
    img = Image.new("RGB", (700, 200), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 26)
    except OSError:
        font = ImageFont.load_default()
    draw.text((20, 20), "REPAIR ESTIMATE", fill="black", font=font)
    draw.text((20, 70), "Total Estimate: $4250.00", fill="black", font=font)
    img.save(str(path))
    return path


@pytest.fixture(scope="module")
def local_http_server(tmp_path_factory):
    """Spins up a real local HTTP server on 127.0.0.1 serving a temp
    directory, so URL ingestion can be tested without external network
    access. Loopback connections aren't subject to outbound network
    restrictions the way real external fetches are.
    """
    serve_dir = tmp_path_factory.mktemp("http_root")

    # Drop a couple of files into the served directory.
    (serve_dir / "claim_note.html").write_text(
        "<html><body><h1>Claim Note</h1>"
        "<p>Policy Number: AUTO-2026-00981</p>"
        "<p>Date of Loss: 2026-06-12</p>"
        "<script>console.log('should be stripped')</script>"
        "</body></html>",
        encoding="utf-8",
    )

    import fitz
    pdf_path = serve_dir / "estimate.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Repair Estimate: $4250.00", fontsize=12)
    doc.save(str(pdf_path))
    doc.close()

    handler = functools.partial(http.server.SimpleHTTPRequestHandler, directory=str(serve_dir))
    server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), handler)
    port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    time.sleep(0.2)  # let the server fully bind before tests fire requests

    yield f"http://127.0.0.1:{port}"

    server.shutdown()


# ---------------------------------------------------------------------------
# PDF tests
# ---------------------------------------------------------------------------

def test_digital_pdf_extracts_text_with_bbox(digital_pdf):
    blocks, page_count, warnings = pdf_parser.parse_pdf(str(digital_pdf))
    assert page_count == 2
    assert len(blocks) >= 2
    assert all(b.extraction_method == "pymupdf_text" for b in blocks)
    assert all(b.bbox is not None for b in blocks)
    assert all(b.confidence == 1.0 for b in blocks)
    full_text = " ".join(b.text for b in blocks)
    assert "AUTO-2026-00981" in full_text


def test_scanned_pdf_triggers_ocr_fallback(scanned_pdf):
    blocks, page_count, warnings = pdf_parser.parse_pdf(str(scanned_pdf))
    assert page_count == 1
    assert len(blocks) > 0
    assert all(b.extraction_method == "pytesseract_ocr" for b in blocks)
    # OCR confidence should be realistic (not hardcoded to 1.0 like digital text)
    assert all(0.0 <= b.confidence <= 1.0 for b in blocks)
    full_text = " ".join(b.text for b in blocks).upper()
    assert "PROP-2026-55210" in full_text.replace(" ", "") or "POLICY" in full_text


def test_long_document_22_pages(long_pdf):
    path, n_pages = long_pdf
    blocks, page_count, warnings = pdf_parser.parse_pdf(str(path))
    assert page_count == n_pages
    assert any("Long document" in w for w in warnings)
    assert len(blocks) >= n_pages  # at least one block per page


def test_pdf_block_ids_are_unique(digital_pdf):
    blocks, _, _ = pdf_parser.parse_pdf(str(digital_pdf))
    ids = [b.block_id for b in blocks]
    assert len(ids) == len(set(ids))


# ---------------------------------------------------------------------------
# DOCX / PPTX tests
# ---------------------------------------------------------------------------

def test_docx_extracts_paragraphs_and_tables(sample_docx):
    blocks, page_count, warnings = office_parser.parse_docx(str(sample_docx))
    assert page_count is None  # honest: docx has no fixed page count
    assert any(b.extraction_method == "docx_paragraph" for b in blocks)
    assert any(b.extraction_method == "docx_table_row" for b in blocks)
    assert all(b.bbox is None for b in blocks)  # flowing doc, no fixed coords
    full_text = " ".join(b.text for b in blocks)
    assert "Daniel Osei" in full_text
    assert "9,840.00" in full_text


def test_pptx_extracts_shapes_with_real_bbox(sample_pptx):
    blocks, page_count, warnings = office_parser.parse_pptx(str(sample_pptx))
    assert page_count == 1
    assert len(blocks) >= 2
    assert all(b.bbox is not None for b in blocks)  # slides ARE a fixed canvas
    assert all(b.page == 1 for b in blocks)
    full_text = " ".join(b.text for b in blocks)
    assert "Priya Nair" in full_text


# ---------------------------------------------------------------------------
# Image (OCR) tests
# ---------------------------------------------------------------------------

def test_standalone_image_ocr(sample_image):
    record = ingest(str(sample_image))
    assert record.source_format == SourceFormat.IMAGE
    assert record.page_count == 1
    assert record.block_count > 0
    full_text = record.full_text.upper()
    assert "ESTIMATE" in full_text


# ---------------------------------------------------------------------------
# Dispatcher-level tests
# ---------------------------------------------------------------------------

def test_dispatcher_routes_by_extension(digital_pdf, sample_docx, sample_pptx, sample_image):
    pdf_rec = ingest(str(digital_pdf))
    docx_rec = ingest(str(sample_docx))
    pptx_rec = ingest(str(sample_pptx))
    img_rec = ingest(str(sample_image))

    assert pdf_rec.source_format == SourceFormat.PDF
    assert docx_rec.source_format == SourceFormat.DOCX
    assert pptx_rec.source_format == SourceFormat.PPTX
    assert img_rec.source_format == SourceFormat.IMAGE


def test_unsupported_extension_raises(tmp_path):
    bad_file = tmp_path / "claim.exe"
    bad_file.write_text("not a real document")
    with pytest.raises(UnsupportedFormatError):
        ingest(str(bad_file))


def test_missing_file_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        ingest(str(tmp_path / "does_not_exist.pdf"))


def test_ingest_many_skips_bad_files_without_aborting_batch(digital_pdf, tmp_path):
    bad_file = tmp_path / "broken.exe"
    bad_file.write_text("garbage")
    records = ingest_many([str(digital_pdf), str(bad_file)])
    assert len(records) == 2
    assert records[0].block_count > 0
    assert records[1].block_count == 0
    assert any("INGESTION FAILED" in w for w in records[1].warnings)


# ---------------------------------------------------------------------------
# Hyperlink ingestion tests (local HTTP server, no external network needed)
# ---------------------------------------------------------------------------

def test_ingest_url_html(local_http_server):
    record = ingest(f"{local_http_server}/claim_note.html")
    assert record.block_count > 0
    full_text = record.full_text
    assert "AUTO-2026-00981" in full_text
    assert "console.log" not in full_text  # <script> content must be stripped


def test_ingest_url_pdf_routes_through_pdf_parser(local_http_server):
    record = ingest(f"{local_http_server}/estimate.pdf")
    assert record.source_format == SourceFormat.PDF
    assert record.block_count > 0
    assert "4250" in record.full_text


def test_ingest_url_404_raises(local_http_server):
    from agents.ingestion.base import FetchError
    with pytest.raises(FetchError):
        ingest(f"{local_http_server}/does_not_exist.html")
