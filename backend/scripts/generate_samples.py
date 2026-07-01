"""
scripts/generate_samples.py
------------------------------
Generates synthetic claim-document fixtures under samples/, one per
supported format, plus a 22-page stress-test PDF (the manager's "20-25
pages" requirement). These are NOT real claim documents -- they're
programmatically generated stand-ins so the ingestion pipeline (and its
tests) have something realistic to chew on without needing real, sensitive
insurance paperwork.

Run: python3 scripts/generate_samples.py
"""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

from core.config import SAMPLES_DIR


def make_auto_claim_pdf(path: Path) -> None:
    """A 2-page digital PDF -- the 'normal' case (real text layer)."""
    doc = fitz.open()
    page = doc.new_page()
    lines = [
        "AUTOMOBILE LOSS NOTICE",
        "",
        "Policy Number: AUTO-2026-00981",
        "Carrier: Meridian Mutual Insurance  NAIC Code: 11223",
        "Policy Period: 2026-01-01 to 2026-12-31",
        "",
        "Date of Loss: 2026-06-12   Time of Loss: 5:45 PM",
        "Location of Loss: Intersection of 5th Ave and Oak St, Springfield",
        "",
        "Description of Accident:",
        "Insured vehicle was stopped at a red light when struck from behind",
        "by the other party's vehicle. Rear bumper and trunk sustained damage.",
        "",
        "Driver Name: Priya Nair",
        "Driver License Number: D1234567",
        "Vehicle VIN: 1HGCM82633A004352",
        "Vehicle Make/Model/Year: Honda Accord 2022",
        "Where Damage Can Be Seen: Rear bumper, trunk lid",
        "",
        "Other Party / Driver Name: Marcus Webb",
        "Other Vehicle Plate: XYZ-4821",
        "",
        "Witness Name: Aaliyah Cho",
    ]
    y = 50
    for line in lines:
        page.insert_text((50, y), line, fontsize=11)
        y += 16

    page2 = doc.new_page()
    lines2 = [
        "AUTOMOBILE LOSS NOTICE (continued)",
        "",
        "Police Report Number: SPD-2026-44210",
        "Reporting Police Department: Springfield Police Department",
        "",
        "Injuries Reported: No",
        "",
        "Repair Estimate: $4,250.00",
        "Rental Vehicle Needed: Yes",
        "",
        "Additional Remarks:",
        "Insured reports no prior accidents in the last 5 years.",
    ]
    y = 50
    for line in lines2:
        page2.insert_text((50, y), line, fontsize=11)
        y += 16

    doc.save(str(path))
    doc.close()


def make_scanned_property_claim_pdf(path: Path) -> None:
    """A 1-page PDF with NO embedded text layer -- an image of text only,
    forcing the OCR fallback path. Simulates a faxed/photographed loss
    notice."""
    img = Image.new("RGB", (1000, 700), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 22)
    except OSError:
        font = ImageFont.load_default()

    text_lines = [
        "PROPERTY LOSS NOTICE",
        "Policy Number: PROP-2026-55210",
        "Date of Loss: 2026-06-10",
        "Cause of Loss: Wind damage to roof shingles",
        "Insured Location Address: 142 Birchwood Lane, Riverton",
        "Probable Amount of Entire Loss: $18,500.00",
        "Inspection Report Reference: INSP-998213",
    ]
    y = 40
    for line in text_lines:
        draw.text((40, y), line, fill="black", font=font)
        y += 45

    doc = fitz.open()
    page = doc.new_page(width=1000, height=700)
    page.insert_image(fitz.Rect(0, 0, 1000, 700), stream=_img_to_png_bytes(img))
    doc.save(str(path))
    doc.close()


def _img_to_png_bytes(img: Image.Image) -> bytes:
    import io
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def make_health_claim_docx(path: Path) -> None:
    doc = DocxDocument()
    doc.add_heading("Health Claim Intake Summary", level=1)
    doc.add_paragraph("Patient Name: Daniel Osei")
    doc.add_paragraph("Member ID: MBR-77231908")
    doc.add_paragraph("Date of Birth: 1989-03-22")
    doc.add_paragraph("Provider / Facility Name: Lakeside General Hospital")
    doc.add_paragraph("Admission Date: 2026-06-08")
    doc.add_paragraph("Discharge Date: 2026-06-11")
    doc.add_paragraph(
        "Diagnosis Description: Acute appendicitis, laparoscopic appendectomy performed."
    )

    doc.add_heading("Billing Summary", level=2)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text, hdr[1].text = "Item", "Amount"
    rows = [
        ("Total Billed Amount", "$9,840.00"),
        ("Co-Pay Amount", "$250.00"),
        ("Allowed Amount", "$8,100.00"),
    ]
    for item, amount in rows:
        row = table.add_row().cells
        row[0].text, row[1].text = item, amount

    doc.add_paragraph("Pre-Authorization Number: PA-2026-30112")
    doc.save(str(path))


def make_claim_overview_pptx(path: Path) -> None:
    """Simulates an internal claim-summary slide deck (plausible real-world
    artifact: an adjuster's case-review deck), to exercise the PPTX parser
    with real shape-level bounding boxes."""
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Claim Overview: AUTO-2026-00981"
    body = slide1.placeholders[1].text_frame
    body.text = "Policyholder: Priya Nair"
    body.add_paragraph().text = "Loss Date: 2026-06-12"
    body.add_paragraph().text = "Status: Pending adjuster review"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Financial Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Repair Estimate: $4,250.00"
    body2.add_paragraph().text = "Rental Coverage: Approved"

    box = slide2.shapes.add_textbox(Inches(1), Inches(4.5), Inches(4), Inches(1))
    box.text_frame.text = "Note: Awaiting police report confirmation."

    prs.save(str(path))


def make_damage_photo_png(path: Path) -> None:
    """Simulates a phone photo of a printed repair estimate -- pure image,
    no text layer, exercises the standalone image OCR path."""
    img = Image.new("RGB", (900, 500), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 26)
    except OSError:
        font = ImageFont.load_default()
    lines = [
        "REPAIR ESTIMATE",
        "Vehicle: Honda Accord 2022",
        "Rear bumper replacement: $1,800.00",
        "Trunk lid repair: $1,450.00",
        "Paint and labor: $1,000.00",
        "Total Estimate: $4,250.00",
    ]
    y = 40
    for line in lines:
        draw.text((40, y), line, fill="black", font=font)
        y += 55
    img.save(str(path))


def make_long_stress_test_pdf(path: Path, n_pages: int = 22) -> None:
    """A 22-page digital PDF, validating the manager's '20-25 pages' requirement
    and the long-document warning threshold."""
    doc = fitz.open()
    for i in range(1, n_pages + 1):
        page = doc.new_page()
        page.insert_text((50, 50), f"CLAIM FILE SUPPLEMENT -- Page {i} of {n_pages}", fontsize=13)
        page.insert_text((50, 80), f"Policy Number: AUTO-2026-00981", fontsize=11)
        page.insert_text(
            (50, 100),
            f"Adjuster note #{i}: Reviewed supporting documentation, no discrepancies found.",
            fontsize=11,
        )
    doc.save(str(path))
    doc.close()


def main() -> None:
    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)

    make_auto_claim_pdf(SAMPLES_DIR / "auto_claim_01.pdf")
    make_scanned_property_claim_pdf(SAMPLES_DIR / "property_claim_01_scanned.pdf")
    make_health_claim_docx(SAMPLES_DIR / "health_claim_01.docx")
    make_claim_overview_pptx(SAMPLES_DIR / "claim_overview_01.pptx")
    make_damage_photo_png(SAMPLES_DIR / "damage_photo_01.png")
    make_long_stress_test_pdf(SAMPLES_DIR / "long_supplement_22pages.pdf", n_pages=22)

    print(f"Generated sample documents in: {SAMPLES_DIR}")
    for f in sorted(SAMPLES_DIR.iterdir()):
        print(f"  - {f.name}")


if __name__ == "__main__":
    main()
